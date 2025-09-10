import streamlit as st
import pandas as pd
import io
import html
import re
import emoji
import os
from google.cloud import translate_v3 as translate
from docx import Document
import PyPDF2
from pptx import Presentation

# Initialize Google Cloud Project and Region
PROJECT_ID = 'ford-180395bd732cdd9af050c1f7'  # Your project ID
REGION = 'us-central1'

# Initialize the Translator client for v3
translate_client_v3 = translate.TranslationServiceClient()

# List of common languages for user selection
# Removed 'Auto Detect' as it's not supported by Translation LLM
LANGUAGES = {
    'Arabic': 'ar',
    'Bulgarian': 'bg',
    'Catalan': 'ca',
    'Chinese (Simplified)': 'zh-CN',
    'Chinese (Traditional)': 'zh-TW',
    'Croatian': 'hr',
    'Czech': 'cs',
    'Danish': 'da',
    'Dutch': 'nl',
    'English': 'en',
    'Estonian': 'et',
    'Finnish': 'fi',
    'French': 'fr',
    'German': 'de',
    'Greek': 'el',
    'Hungarian': 'hu',
    'Indonesian': 'id',
    'Italian': 'it',
    'Japanese': 'ja',
    'Korean': 'ko',
    'Latvian': 'lv',
    'Lithuanian': 'lt',
    'Polish': 'pl',
    'Portuguese': 'pt',
    'Romanian': 'ro',
    'Russian': 'ru',
    'Slovak': 'sk',
    'Slovenian': 'sl',
    'Spanish': 'es',
    'Swedish': 'sv',
    'Thai': 'th',
    'Turkish': 'tr',
    'Ukrainian': 'uk',
    'Vietnamese': 'vi'
}

# --- Text Pre/Post-processing Functions (kept as is) ---
def convert_emoticons(text):
    return emoji.demojize(text)

def remove_usernames(text):
    return re.sub(r'@\w+', '', text)

def remove_hyperlinks(text):
    return re.sub(r'http\S+|www\S+|https\S+', '', text)

def clean_extra_spaces(text):
    return re.sub(r'\s+', ' ', text).strip()

def standardize_quotes(text):
    return text.replace('“', '"').replace('”', '"')

def remove_special_characters(text):
    return re.sub(r'[^a-zA-Z0-9\s.,!?\'"()]+', '', text)

# --- Updated Translation Function using translate_v3 and Translation LLM ---
# --- Updated Translation Function using translate_v3 and Translation LLM ---
def translate_text_with_llm(text, source_language_code, project_id, region):
    """
    Translates text to English using the Translation LLM model (v3 API).
    Includes post-processing steps.
    Requires an explicit source_language_code.
    """
    if not text or not isinstance(text, str):
        return "" # Return empty string for non-string or empty inputs

    if not source_language_code:
        return "Error: Source language must be explicitly selected for Translation LLM."

    parent = f"projects/{project_id}/locations/{region}"
    
    # *** CRITICAL CHANGE HERE ***
    # Reference the Translation LLM as a 'general' model, not a project-specific one.
    model_path = f"projects/{project_id}/locations/{region}/models/general/translation-llm"

    try:
        request_body = {
            "parent": parent,
            "contents": [text],
            "target_language_code": 'en',
            "source_language_code": source_language_code, # Explicitly required
            "model": model_path,
        }

        response = translate_client_v3.translate_text(request=request_body)
        translated_text = response.translations[0].translated_text
        
        # Apply post-processing
        translated_text = html.unescape(translated_text)
        translated_text = convert_emoticons(translated_text)
        translated_text = remove_usernames(translated_text)
        translated_text = remove_hyperlinks(translated_text)
        translated_text = clean_extra_spaces(translated_text)
        translated_text = standardize_quotes(translated_text)
        translated_text = remove_special_characters(translated_text)

        return translated_text
    except Exception as e:
        st.error(f"Error during LLM translation: {str(e)}")
        return f"Translation Error: {str(e)}"

# Function to extract text from various file types (kept as is)
def extract_text_from_file(file):
    if file.name.endswith('.xlsx'):
        return pd.read_excel(file, sheet_name=None)
    elif file.name.endswith('.csv'):
        file_content = file.read().decode('utf-8')
        return {'Sheet1': pd.read_csv(io.StringIO(file_content))}
    elif file.name.endswith('.docx'):
        doc = Document(file)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file.name.endswith('.pdf'):
        pdf_reader = PyPDF2.PdfReader(file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text() + '\n'
        return text
    elif file.name.endswith('.pptx'):
        ppt = Presentation(file)
        text = ''
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    else:
        return None

# Streamlit Application
st.title("Multi-Document Language Translation App (Powered by Translation LLM and Google Translate V3)")
st.write("Upload a document with text for translation.")
st.write("This application uses Google Cloud Translation v3 with the Translation LLM for text content.")
st.warning("Note: The Translation LLM requires you to explicitly select the source language. 'Auto Detect' is not supported.")

# Language selection dropdown (now without 'Auto Detect')
source_language_display = st.selectbox("Select source language:", list(LANGUAGES.keys()))
source_language_code = LANGUAGES[source_language_display]

# File uploader for multiple file types
uploaded_file = st.file_uploader("Choose a file", type=['xlsx', 'csv', 'docx', 'pdf', 'pptx'])

if uploaded_file is not None:
    extracted_content = extract_text_from_file(uploaded_file)

    if isinstance(extracted_content, dict):
        st.write("Extracted sheets/dataframes:")
        processed_data = {}

        for sheet_name, df in extracted_content.items():
            st.write(f"**{sheet_name} (Original)**")
            st.dataframe(df.astype(str)) # Convert to string for display to avoid pyarrow issues
            processed_data[sheet_name] = df.copy()

        if st.button("Translate Excel/CSV"):
            with st.spinner("Translating... This may take a while for large files."):
                for sheet_name, df_to_process in processed_data.items():
                    for column in df_to_process.columns:
                        df_to_process[column] = df_to_process[column].astype(str).apply(
                            lambda x: translate_text_with_llm(x, source_language_code, PROJECT_ID, REGION)
                        )
                    processed_data[sheet_name] = df_to_process

            st.success("Translation complete!")

            for sheet_name, df_translated in processed_data.items():
                st.write(f"**{sheet_name} (Translated)**")
                st.dataframe(df_translated.astype(str)) # Convert to string for display

            output_buffer = io.BytesIO()
            if uploaded_file.name.endswith('.xlsx'):
                with pd.ExcelWriter(output_buffer, engine='openpyxl') as writer:
                    for sheet_name, data in processed_data.items():
                        data.to_excel(writer, index=False, sheet_name=sheet_name)
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                original_file_name = os.path.splitext(uploaded_file.name)[0]
                output_file_name = f'translated_{original_file_name}.xlsx'
            elif uploaded_file.name.endswith('.csv'):
                df_translated = list(processed_data.values())[0]
                df_translated.to_csv(output_buffer, index=False, encoding='utf-8')
                mime_type = 'text/csv'
                original_file_name = os.path.splitext(uploaded_file.name)[0]
                output_file_name = f'translated_{original_file_name}.csv'

            output_buffer.seek(0)

            st.download_button(
                label="Download Translated File",
                data=output_buffer,
                file_name=output_file_name,
                mime=mime_type
            )

    elif isinstance(extracted_content, str):
        st.write("Extracted Text:")
        st.text_area("Original Text", extracted_content, height=300)

        if st.button("Translate Text"):
            with st.spinner("Translating..."):
                translated_text = translate_text_with_llm(extracted_content, source_language_code, PROJECT_ID, REGION)
            
            st.success("Translation complete!")
            st.write("Translated Text:")
            st.text_area("Translated Text", translated_text, height=300)

            output_buffer = io.BytesIO()
            output_buffer.write(translated_text.encode('utf-8'))
            output_buffer.seek(0)
            
            st.download_button(
                label="Download Translated Text",
                data=output_buffer,
                file_name='translated_text.txt',
                mime='text/plain'
            )
    
    else:
        st.error("Unsupported file type or no text extracted from the file.")
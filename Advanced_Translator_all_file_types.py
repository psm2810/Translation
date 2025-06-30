import streamlit as st
import pandas as pd
import io
import html
import re
import emoji
import os
from google.cloud import aiplatform
from google.cloud import translate_v2 as translate
from docx import Document
import PyPDF2
from pptx import Presentation

# Initialize Google Cloud AI Platform
PROJECT_ID = 'ford-180395bd732cdd9af050c1f7'  # Your project ID
REGION = 'us-central1'  # Your region
aiplatform.init(project=PROJECT_ID, location=REGION)

# Initialize the Translator client
translate_client = translate.Client()

# List of common languages for user selection
LANGUAGES = {
    'Auto Detect': None,
    'Arabic': 'ar',
    'Bulgarian': 'bg',
    'Catalan': 'ca',
    'Chinese (Simplified)': 'zh-CN',
    'Chinese (Traditional)': 'zh-TW',
    'Croatian': 'hr',
    'Czech': 'cs',
    'Danish': 'da',
    'Dutch': 'nl',
    'English': 'en',
    'Estonian': 'et',
    'Finnish': 'fi',
    'French': 'fr',
    'German': 'de',
    'Greek': 'el',
    'Hungarian': 'hu',
    'Indonesian': 'id',
    'Italian': 'it',
    'Japanese': 'ja',
    'Korean': 'ko',
    'Latvian': 'lv',
    'Lithuanian': 'lt',
    'Polish': 'pl',
    'Portuguese': 'pt',
    'Romanian': 'ro',
    'Russian': 'ru',
    'Slovak': 'sk',
    'Slovenian': 'sl',
    'Spanish': 'es',
    'Swedish': 'sv',
    'Thai': 'th',
    'Turkish': 'tr',
    'Ukrainian': 'uk',
    'Vietnamese': 'vi'
}

def convert_emoticons(text):
    return emoji.demojize(text)

def remove_usernames(text):
    return re.sub(r'@\w+', '', text)

def remove_hyperlinks(text):
    return re.sub(r'http\S+|www\S+|https\S+', '', text)

def clean_extra_spaces(text):
    return re.sub(r'\s+', ' ', text).strip()

def standardize_quotes(text):
    return text.replace('“', '"').replace('”', '"')

def remove_special_characters(text):
    return re.sub(r'[^a-zA-Z0-9\s.,!?\'"()]+', '', text)

# Function to translate text to English with post-processing
def translate_text(text, source_language):
    try:
        result = translate_client.translate(text, target_language='en', source_language=source_language)
        translated_text = html.unescape(result['translatedText'])
        
        # Apply post-processing
        translated_text = convert_emoticons(translated_text)
        translated_text = remove_usernames(translated_text)
        translated_text = remove_hyperlinks(translated_text)
        translated_text = clean_extra_spaces(translated_text)
        translated_text = standardize_quotes(translated_text)
        translated_text = remove_special_characters(translated_text)

        return translated_text
    except Exception as e:
        return f"Error: {str(e)}"

# Function to extract text from various file types
def extract_text_from_file(file):
    if file.name.endswith('.xlsx'):
        return pd.read_excel(file, sheet_name=None)  # Read all sheets into a dict
    elif file.name.endswith('.csv'):
        return pd.read_csv(file)  # Read CSV
    elif file.name.endswith('.docx'):
        doc = Document(file)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file.name.endswith('.pdf'):
        pdf_reader = PyPDF2.PdfReader(file)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text() + '\n'
        return text
    elif file.name.endswith('.pptx'):
        ppt = Presentation(file)
        text = ''
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text
    else:
        return None

# Streamlit Application
st.title("Multi-Document Language Translation App")
st.write("Upload a document with text for translation.")

# Language selection dropdown
source_language = st.selectbox("Select source language:", list(LANGUAGES.keys()))

# File uploader for multiple file types
uploaded_file = st.file_uploader("Choose a file", type=['xlsx', 'csv', 'docx', 'pdf', 'pptx'])

if uploaded_file is not None:
    extracted_text = extract_text_from_file(uploaded_file)

    # Handle different cases based on file type
    if isinstance(extracted_text, dict):  # If it's an Excel file with multiple sheets
        st.write("Extracted sheets:")
        sheet_data = {}
        for sheet_name, df in extracted_text.items():
            st.write(f"**{sheet_name}**")
            st.dataframe(df)
            sheet_data[sheet_name] = df  # Store original data for processing later

        # Submit button to trigger translation
        if st.button("Submit"):
            # Translate all sheets
            for sheet_name, df in sheet_data.items():
                for column in df.columns:
                    df[column] = df[column].astype(str).apply(lambda x: translate_text(x, LANGUAGES[source_language]))
                sheet_data[sheet_name] = df  # Update the translated DataFrame
            
            # Create a BytesIO buffer for the Excel file
            output_buffer = io.BytesIO()
            with pd.ExcelWriter(output_buffer, engine='openpyxl') as writer:
                for sheet_name, data in sheet_data.items():
                    data.to_excel(writer, index=False, sheet_name=sheet_name)

            # Set the buffer position to the beginning
            output_buffer.seek(0)

            # Construct the output file name
            original_file_name = os.path.splitext(uploaded_file.name)[0]  # Get the original file name without extension
            output_file_name = f'translated_{original_file_name}.xlsx'  # Add the translated prefix

            # Download button for the translated output
            st.download_button(
                label="Download Translated File",
                data=output_buffer,
                file_name=output_file_name,
                mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            )

    elif isinstance(extracted_text, str):  # If it's a text extracted from docx, pdf, or pptx
        st.write("Extracted Text:")
        st.text_area("Text for Translation", extracted_text, height=300)

        # Submit button to trigger translation
        if st.button("Submit"):
            # Translate the text
            translated_text = translate_text(extracted_text, LANGUAGES[source_language])
            st.write("Translated Text:")
            st.text_area("Translated Text", translated_text, height=300)

            # Download button for the translated output
            if st.button("Download Translated Text"):
                # Create a text file for download
                output_buffer = io.BytesIO()
                output_buffer.write(translated_text.encode('utf-8'))
                output_buffer.seek(0)
                
                st.download_button(
                    label="Download Translated File",
                    data=output_buffer,
                    file_name='translated_text.txt',
                    mime='text/plain'
                )
    
    else:
        st.error("Unsupported file type or empty file.")
